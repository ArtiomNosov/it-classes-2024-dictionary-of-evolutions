import base64
import os
import sqlite3
import requests
import io
import numpy as np
from flask import Flask, render_template, request, jsonify, send_file
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from bs4 import BeautifulSoup
app = Flask(__name__, static_folder="static", template_folder="templates")


def gpt_response(subject, century, number):
    url = "http://localhost:1337/v1/chat/completions"
    body = {
        "model": "gpt-3.5-turbo",
        "stream": False,
        "messages": [
            {"role": "assistant",
             "content": f"Расскажи про {subject} в {century} веке для {number} слайда презентации, уложись в 30 слов полностью на русском, скинь только то что надо вставить в презентацию"}
        ]
    }

    json_response = requests.post(url, json=body).json().get('choices', [])
    g = []
    for choice in json_response:
        g.append(choice.get('message', {}).get('content', ''))
    return " ".join(g)

def google_search(query):
    url = f"https://www.google.com/search?q={query}"
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        soup = BeautifulSoup(response.text, "html.parser")
        first_result = soup.select_one(".tF2Cxc")
        if first_result:
            title = first_result.select_one(".DKV0Md").text
            link = first_result.select_one(".yuRUbf a")["href"]
            snippet = first_result.select_one(".VwiC3b").text
            return {"title": title, "link": link, "snippet": snippet}
        else:
            return "Результатов не найдено"
    else:
        return f"Ошибка при выполнении запроса: {response.status_code}"
def remove_bg(image_path, threshold=240):
    image = Image.open(image_path).convert("RGBA")
    data = np.array(image)
    mask = (data[:, :, :3].sum(axis=2) > threshold * 3)
    data[mask, 3] = 0
    result_image = Image.fromarray(data, mode="RGBA")
    filename, file_extension = os.path.splitext(image_path)
    result_image.save(filename + '.png')


def add_resized_picture(slide, img_path, max_width, max_height, left, top):
    img = Image.open(img_path)
    width, height = img.size

    scale = min(max_width / width, max_height / height)
    new_width = int(width * scale)
    new_height = int(height * scale)

    slide.shapes.add_picture(img_path, left, top, width=Inches(new_width / 96), height=Inches(new_height / 96))


def search_exact_items_as_dict(item_name):
    conn = sqlite3.connect('items.db')
    cursor = conn.cursor()

    cursor.execute("SELECT century, item_name, image FROM items WHERE item_name = ?", (item_name,))
    rows = cursor.fetchall()

    items = []
    for row in rows:
        item_dict = {
            "century": row[0],
            "item_name": row[1],
            "image": row[2]
        }
        items.append(item_dict)

    conn.close()
    return items


def int_to_roman(n):
    roman_numerals = {
        1: "I", 4: "IV", 5: "V", 9: "IX", 10: "X",
        40: "XL", 50: "L", 90: "XC", 100: "C",
        400: "CD", 500: "D", 900: "CM", 1000: "M"
    }
    result = ""
    for value in sorted(roman_numerals.keys(), reverse=True):
        while n >= value:
            result += roman_numerals[value]
            n -= value
    return result


def create_of_presentation():
    a = input()
    prs = Presentation()
    search_query = a
    items = sorted(search_exact_items_as_dict(search_query), key=lambda x: int(x["century"]))

    c = 0
    for i in items:
        c += 1
        century = i["century"]
        sub = i["item_name"]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        img_path = 'img1/back.jpg'
        slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)
        add_resized_picture(slide, "img1/img.png", 1 * 96, 1 * 96, left=Inches(9), top=Inches(6.5))
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title = title_box.text_frame
        title.text = f"{sub} {century} век"
        title.paragraphs[0].font.size = Pt(36)
        title.paragraphs[0].font.bold = True
        title.paragraphs[0].font.name = "Times New Roman"
        title.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4), Inches(8))
        text = text_box.text_frame
        f = gpt_response(sub, century, c)
        while ("Извините" in f) or (f == "") or (
                f == " ") or f == "Model not found or too long input. Or any other error (xD)" or (
                "a" in f or "b" in f or "c" in f or "d" in f or "e" in f or
                "f" in f or "g" in f or "h" in f or "i" in f or "j" in f or
                "k" in f or "l" in f or "m" in f or "n" in f or "o" in f or
                "p" in f or "r" in f or "s" in f or "t" in f or "u" in f or
                "v" in f or "x" in f or "y" in f or "z" in f or "q" in f or
                "w" in f or
                "A" in f or "B" in f or "C" in f or "D" in f or "E" in f or
                "F" in f or "G" in f or "H" in f or "I" in f or "J" in f or
                "K" in f or "L" in f or "M" in f or "N" in f or "O" in f or
                "P" in f or "R" in f or "S" in f or "T" in f or "U" in f or
                "V" in f or "X" in f or "Y" in f or "Z" in f
        ) or "Удачи" in f or "удачи" in f or "слайд" in f or "Слайд" in f or "позитив" in f:
            f = gpt_response(sub, century, c)
        text.text = f
        text.paragraphs[0].font.size = Pt(18)
        text.paragraphs[0].font.name = "Times New Roman"
        text.word_wrap = True
        image = Image.open(io.BytesIO(i["image"]))
        image_path = f"{sub}.png"
        image.save(image_path)
        remove_bg(image_path, 235)
        max_width = 6 * 96
        max_height = 3 * 96
        add_resized_picture(slide, "img1/icon.png", 1 * 56, 1 * 56, left=Inches(4.5), top=Inches(1.95))
        add_resized_picture(slide, image_path, max_width, max_height, left=Inches(1.2), top=Inches(2))
        os.remove(image_path)

        start_x = Inches(1)
        y_position = Inches(5.5)
        line_length = Inches(8)
        num_centuries = 5
        century_spacing = line_length / (num_centuries - 1)
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, y_position, start_x + line_length, y_position
        )
        line_shape.line.width = Pt(2)
        line_shape.line.color.rgb = RGBColor(0, 0, 0)

        centuries = [
            (f"{int_to_roman(int(century) - 2)} век"),
            (f"{int_to_roman(int(century) - 1)} век"),
            (f"{int_to_roman(int(century))} век"),
            (f"{int_to_roman(int(century) + 1)} век"),
            (f"{int_to_roman(int(century) + 2)} век")
        ]

        for i, century in enumerate(centuries):
            x_position = start_x + i * century_spacing

            century_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left=x_position,
                top=y_position - Inches(0.1),
                width=Pt(12),
                height=Pt(12)
            )
            century_shape.fill.solid()
            if i == 0 or i == 1 or i == 3 or i == 4:
                century_shape.fill.fore_color.rgb = RGBColor(0, 0, 255)
            else:
                century_shape.fill.fore_color.rgb = RGBColor(220, 20, 60)
            textbox = slide.shapes.add_textbox(
                left=x_position - Inches(0.3),
                top=y_position + Inches(0.3),
                width=Inches(1),
                height=Inches(0.5)
            )
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = century
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    prs.save('1.pptx')


@app.route('/')
def main_page():
    return render_template("main-page.html")


@app.route('/templates/main-page.html')
def main_page1():
    return render_template("main-page.html")


@app.route('/templates/authors.html')
def authors_page():
    return render_template("authors.html")


@app.route('/templates/review.html')
def review_page():
    return render_template("review.html")


@app.route('/templates/login.html')
def login_page():
    return render_template("login.html")


@app.route('/templates/sign-up.html')
def sign_up_page():
    return render_template("sign-up.html")


@app.route('/templates/pictures.html', methods=['GET'])
def pictures_page():
    item_name = request.args.get('item_name', '').lower().replace(' ', '')
    conn = sqlite3.connect('items.db')
    cursor = conn.cursor()

    # Запрос данных из базы с сортировкой по веку
    cursor.execute("SELECT id, item_name, image, century FROM items WHERE item_name LIKE ? ORDER BY century ASC",
                   ('%' + item_name + '%',))
    rows = cursor.fetchall()
    conn.close()

    images = []
    for row in rows:
        image_data = row[2]
        image_id = row[0]
        if image_data:
            print(f"Image found for item {row[1]}, ID: {image_id}, size: {len(image_data)} bytes")
        else:
            print(f"No image found for item {row[1]} with ID: {image_id}")
        encoded_image_data = base64.b64encode(image_data).decode('utf-8')
        images.append({
            'id': image_id,
            'image_data': encoded_image_data,
            'century': row[3]
        })

    # Получение дополнительной информации из Google
    try:
        query = f"Эволюция {item_name}"
        first_result = google_search(query)
        if isinstance(first_result, dict):
            item_text = first_result["snippet"]
            item_link = first_result["link"]
        else:
            item_text = first_result
            item_link = "https://www.google.ru/?hl=ru"
    except:
        item_text = item_name
        item_link = "https://www.google.ru/?hl=ru"

    # Передача данных в шаблон
    return render_template('pictures.html', item_name=item_name.title(), images=images, item_text=item_text,
                           item_link=item_link)


if __name__ == "__main__":
    app.run(debug=True, port=8888)
